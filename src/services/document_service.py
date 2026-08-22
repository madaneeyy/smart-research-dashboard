"""Document ingestion service for Smart Research AI.

This service is intentionally limited to extraction/normalization. Chunking
and retrieval are handled by the dedicated uploaded-document pipeline.
"""
from __future__ import annotations

import csv
import io
import json
import re
from pathlib import Path
from typing import Any, Dict, Optional
from uuid import uuid4


class DocumentService:
    MAX_FILE_BYTES = 25 * 1024 * 1024
    MAX_TEXT_CHARS = 2_000_000

    TEXT_EXTENSIONS = {
        ".txt", ".md", ".markdown", ".rst",
        ".py", ".pyi", ".js", ".jsx", ".mjs", ".cjs",
        ".ts", ".tsx", ".java", ".c", ".h", ".cc", ".cpp", ".hpp",
        ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".kts", ".dart", ".cs",
        ".scala", ".r", ".jl", ".lua", ".sh", ".bash", ".bat", ".ps1",
        ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".xml",
        ".html", ".htm", ".css", ".sql", ".env", ".log", ".gitignore",
    }

    @classmethod
    def _limit_text(cls, text: str) -> str:
        text = text.replace("\x00", "")
        text = re.sub(r"\n{4,}", "\n\n\n", text)
        text = text.strip()
        if len(text) > cls.MAX_TEXT_CHARS:
            text = text[: cls.MAX_TEXT_CHARS]
            text = text.rsplit("\n", 1)[0]
        return text

    @classmethod
    def extract(
        cls,
        filename: str,
        raw_bytes: bytes,
        content_type: Optional[str] = None,
    ) -> Dict[str, Any]:
        if not filename:
            raise ValueError("Uploaded file has no filename.")
        if len(raw_bytes) > cls.MAX_FILE_BYTES:
            raise ValueError(
                f"File is too large. Maximum size is {cls.MAX_FILE_BYTES // (1024 * 1024)} MB."
            )

        suffix = Path(filename).suffix.lower()

        if suffix == ".pdf":
            text, pages = cls._extract_pdf(raw_bytes)
        elif suffix == ".docx":
            text, pages = cls._extract_docx(raw_bytes)
        elif suffix == ".csv":
            text, pages = cls._extract_csv(raw_bytes)
        elif suffix == ".json":
            text, pages = cls._extract_json(raw_bytes)
        elif suffix in {".xlsx", ".xls"}:
            text, pages = cls._extract_spreadsheet(raw_bytes, suffix)
        elif suffix == ".pptx":
            text, pages = cls._extract_pptx(raw_bytes)
        elif suffix in cls.TEXT_EXTENSIONS or not suffix:
            text = cls._decode_text(raw_bytes)
            pages = None
        else:
            raise ValueError(
                f"Unsupported file type '{suffix or 'unknown'}'. Supported types include "
                "PDF, DOCX, TXT, Markdown, CSV, JSON, XLSX, PPTX and common code/text files."
            )

        text = cls._limit_text(text)
        if not text:
            raise ValueError(
                "The uploaded document contains no extractable text. "
                "If this is a scanned PDF, install OCR support and try again."
            )

        return {
            "document_id": uuid4().hex,
            "filename": filename,
            "path": filename,
            "content": text,
            "source": "upload",
            "category": "document",
            "language": suffix.lstrip(".") or "text",
            "content_type": content_type or "",
            "size_bytes": len(raw_bytes),
            "pages": pages,
            "characters": len(text),
        }

    @staticmethod
    def _decode_text(raw_bytes: bytes) -> str:
        for encoding in ("utf-8", "utf-8-sig", "utf-16", "latin-1"):
            try:
                return raw_bytes.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise ValueError("Could not decode the text file.")

    @classmethod
    def _extract_pdf(cls, raw_bytes: bytes) -> tuple[str, int]:
        # First choice: pypdf, because it is lightweight.
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw_bytes))
            parts = []
            for index, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    parts.append(f"[Page {index}]\n{page_text.strip()}")
            text = "\n\n".join(parts)
            if text.strip():
                return text, len(reader.pages)
        except ImportError:
            pass
        except Exception:
            # Fall through to PyMuPDF where available.
            pass

        # PyMuPDF is often better at real-world PDFs with complicated layouts.
        try:
            import fitz
            pdf = fitz.open(stream=raw_bytes, filetype="pdf")
            parts = []
            for index, page in enumerate(pdf, start=1):
                page_text = page.get_text("text") or ""
                if page_text.strip():
                    parts.append(f"[Page {index}]\n{page_text.strip()}")
            text = "\n\n".join(parts)
            if text.strip():
                return text, len(pdf)
        except ImportError:
            pass
        except Exception:
            pass

        raise ValueError(
            "Could not extract text from this PDF. The PDF may be scanned/image-only, "
            "encrypted, or require OCR. Install pypdf or pymupdf; OCR can be added for "
            "image-only PDFs."
        )

    @classmethod
    def _extract_docx(cls, raw_bytes: bytes) -> tuple[str, int]:
        try:
            from docx import Document
        except ImportError as exc:
            raise RuntimeError(
                "DOCX support requires 'python-docx'. Install it with: pip install python-docx"
            ) from exc

        document = Document(io.BytesIO(raw_bytes))
        parts = []
        for paragraph in document.paragraphs:
            value = paragraph.text.strip()
            if not value:
                continue
            style = str(getattr(paragraph.style, "name", "") or "")
            if style.lower().startswith("heading"):
                parts.append(f"# {value}")
            else:
                parts.append(value)

        for table_index, table in enumerate(document.tables, start=1):
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            rows = [r for r in rows if r.strip()]
            if rows:
                parts.append(f"[Table {table_index}]\n" + "\n".join(rows))

        return "\n\n".join(parts), 1

    @classmethod
    def _extract_csv(cls, raw_bytes: bytes) -> tuple[str, int]:
        text = cls._decode_text(raw_bytes)
        reader = csv.reader(io.StringIO(text))
        rows = [" | ".join(str(cell).strip() for cell in row) for row in reader]
        return "\n".join(rows), 1

    @classmethod
    def _extract_json(cls, raw_bytes: bytes) -> tuple[str, int]:
        text = cls._decode_text(raw_bytes)
        try:
            value = json.loads(text)
        except json.JSONDecodeError:
            return text, 1
        return json.dumps(value, indent=2, ensure_ascii=False), 1

    @classmethod
    def _extract_spreadsheet(cls, raw_bytes: bytes, suffix: str) -> tuple[str, int]:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise RuntimeError(
                "Spreadsheet support requires 'openpyxl'. Install it with: pip install openpyxl"
            ) from exc
        if suffix == ".xls":
            raise ValueError("Legacy .xls files are not supported directly; save as .xlsx or .csv.")
        workbook = load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
        parts = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(v.strip() for v in values):
                    rows.append(" | ".join(values))
            if rows:
                parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        return "\n\n".join(parts), len(workbook.worksheets)

    @classmethod
    def _extract_pptx(cls, raw_bytes: bytes) -> tuple[str, int]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError(
                "PPTX support requires 'python-pptx'. Install it with: pip install python-pptx"
            ) from exc
        presentation = Presentation(io.BytesIO(raw_bytes))
        parts = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text).strip():
                    slide_parts.append(str(shape.text).strip())
            if slide_parts:
                parts.append(f"[Slide {index}]\n" + "\n".join(slide_parts))
        return "\n\n".join(parts), len(presentation.slides)
